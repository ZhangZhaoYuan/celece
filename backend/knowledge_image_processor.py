"""
方案C - 知识文档图文联合索引系统
集成到小赛助手主系统

功能：
1. 从Word/PPT文档中提取内嵌图片
2. 使用视觉模型池生成图片描述（自动故障切换）
3. 生成向量并存储到sqlite-vec
4. 建立图文关联，支持双向检索
"""

import os
import re
import json
import hashlib
import base64
import io
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Tuple
import requests

# 配置路径
BASE_DIR = Path(__file__).resolve().parent.parent
DATA_DIR = BASE_DIR / "data"
KNOWLEDGE_IMAGES_DIR = DATA_DIR / "knowledge_images"
IMAGE_VECTORS_DB = DATA_DIR / "image_vectors.db"

# 确保目录存在
KNOWLEDGE_IMAGES_DIR.mkdir(parents=True, exist_ok=True)

# ==================== 模型池配置 ====================

# 视觉模型池（按优先级排序）
VISION_MODEL_POOL = [
    "qwen3.5-omni-plus",
    "qwen3.5-omni-flash",
    "qwen3-omni-flash",
    "qwen-omni-turbo",
    "qwen3-vl-plus",
    "qwen3-vl-flash",
    "qwen-vl-max",
    "qwen-vl-plus",
]

# 向量模型池
EMBEDDING_MODEL_POOL = [
    "text-embedding-v4",
    "text-embedding-v3",
    "text-embedding-v2",
    "qwen3.7-text-embedding",
]

# 模型指针
_vision_model_idx = 0
_embedding_model_idx = 0


def get_next_vision_model() -> str:
    """获取下一个视觉模型（循环）"""
    global _vision_model_idx
    model = VISION_MODEL_POOL[_vision_model_idx % len(VISION_MODEL_POOL)]
    _vision_model_idx += 1
    return model


def get_next_embedding_model() -> str:
    """获取下一个向量模型（循环）"""
    global _embedding_model_idx
    model = EMBEDDING_MODEL_POOL[_embedding_model_idx % len(EMBEDDING_MODEL_POOL)]
    _embedding_model_idx += 1
    return model


# ==================== API配置 ====================

def get_api_config() -> Tuple[str, str]:
    """获取API配置"""
    try:
        import config_manager as cfg
        v_cfg = cfg.get_vision_model_config()
        if v_cfg and v_cfg.get("api_key"):
            return v_cfg["api_key"], v_cfg["base_url"].rstrip("/")
    except:
        pass
    
    # 默认配置
    return "", ""


# ==================== 图片提取 ====================

def extract_images_from_docx(filepath: str) -> List[Dict]:
    """从docx文件中提取所有内嵌图片"""
    try:
        from docx import Document
        
        doc = Document(filepath)
        images = []
        
        for relation in doc.part.rels.values():
            if "image" in relation.reltype:
                try:
                    image = relation.target_part
                    image_data = image.blob
                    ext = image.content_type.split('/')[-1]
                    
                    img_hash = hashlib.md5(image_data).hexdigest()[:12]
                    filename = f"doc_{img_hash}.{ext}"
                    save_path = KNOWLEDGE_IMAGES_DIR / filename
                    
                    if not save_path.exists():
                        save_path.write_bytes(image_data)
                    
                    images.append({
                        'hash': img_hash,
                        'filename': filename,
                        'filepath': str(save_path),
                        'content_type': image.content_type,
                        'size': len(image_data)
                    })
                except Exception as e:
                    print(f"    ⚠️  处理图片失败: {e}")
        
        return images
        
    except ImportError:
        print("  ❌ 需要安装 python-docx: pip install python-docx")
        return []
    except Exception as e:
        print(f"  ❌ 提取图片失败: {e}")
        return []


def extract_images_from_pptx(filepath: str) -> List[Dict]:
    """从pptx文件中提取图片"""
    try:
        from pptx import Presentation
        
        prs = Presentation(filepath)
        images = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        image_data = image.blob
                        ext = image.content_type.split('/')[-1]
                        
                        img_hash = hashlib.md5(image_data).hexdigest()[:12]
                        filename = f"slide{slide_num}_{img_hash}.{ext}"
                        save_path = KNOWLEDGE_IMAGES_DIR / filename
                        
                        if not save_path.exists():
                            save_path.write_bytes(image_data)
                        
                        images.append({
                            'hash': img_hash,
                            'filename': filename,
                            'filepath': str(save_path),
                            'content_type': image.content_type,
                            'size': len(image_data)
                        })
                    except Exception as e:
                        print(f"    ⚠️  处理幻灯片图片失败: {e}")
        
        return images
        
    except ImportError:
        print("  ❌ 需要安装 python-pptx: pip install python-pptx")
        return []
    except Exception as e:
        print(f"  ❌ 提取图片失败: {e}")
        return []


# ==================== 视觉模型调用 ====================

def describe_image(image_path: str, prompt: str = None) -> Tuple[str, str]:
    """
    使用视觉模型池描述图片，自动故障切换
    返回: (描述文本, 使用的模型名)
    """
    api_key, base_url = get_api_config()
    
    if not api_key:
        return "", "未配置API Key"
    
    # 读取图片
    try:
        with open(image_path, 'rb') as f:
            image_data = f.read()
        base64_image = base64.b64encode(image_data).decode()
        ext = Path(image_path).suffix.lower()
        mime_type = 'image/jpeg' if ext in ['.jpg', '.jpeg'] else 'image/png'
    except Exception as e:
        return "", f"读取图片失败: {e}"
    
    # 默认提示词
    if not prompt:
        prompt = """请详细分析这张图片，输出以下信息：

【图片类型】
判断图片类型：排油排便图 / 身材对比照 / 付款截图 / 产品图 / 聊天记录 / 其他

【内容描述】
详细描述图片内容

【适用场景】
这张图片适合在什么销售场景中使用

请按格式输出，每个部分用空行分隔。"""
    
    # 尝试每个模型
    for _ in range(len(VISION_MODEL_POOL)):
        model = get_next_vision_model()
        try:
            resp = requests.post(
                f"{base_url.rstrip('/')}/chat/completions",
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": model,
                    "messages": [{
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{base64_image}"}}
                        ]
                    }],
                    "max_tokens": 500
                },
                timeout=60
            )
            
            if resp.status_code == 200:
                data = resp.json()
                description = data['choices'][0]['message']['content'].strip()
                return description, model
            else:
                print(f"    ⚠️  模型 {model} 失败: {resp.status_code}")
                
        except Exception as e:
            print(f"    ⚠️  模型 {model} 错误: {e}")
    
    return "图片描述生成失败", "全部失败"


# ==================== 向量生成 ====================

def generate_embedding(text: str) -> Tuple[List[float], str]:
    """
    为文本生成向量
    返回: (向量列表, 使用的模型名)
    """
    api_key, base_url = get_api_config()
    
    if not api_key:
        return [], "未配置API Key"
    
    # 尝试每个模型
    for _ in range(len(EMBEDDING_MODEL_POOL)):
        model = get_next_embedding_model()
        try:
            resp = requests.post(
                f"{base_url.rstrip('/')}/embeddings",
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                },
                json={
                    "model": model,
                    "input": text
                },
                timeout=30
            )
            
            if resp.status_code == 200:
                data = resp.json()
                vector = data['data'][0]['embedding']
                return vector, model
            else:
                print(f"    ⚠️  向量模型 {model} 失败: {resp.status_code}")
                
        except Exception as e:
            print(f"    ⚠️  向量模型 {model} 错误: {e}")
    
    return [], "全部失败"


# ==================== 数据库操作 ====================

def init_document_image_table():
    """初始化知识文档图片表"""
    import sqlite3
    
    conn = sqlite3.connect(str(IMAGE_VECTORS_DB))
    
    # 尝试加载sqlite-vec扩展
    try:
        conn.enable_load_extension(True)
        
        # 寻找vec0.dll的路径
        vec_dll_paths = [
            os.path.join(os.path.dirname(__file__), 'sqlite_vec', 'vec0.dll'),
            os.path.join(BASE_DIR, 'sqlite_vec', 'vec0.dll'),
        ]
        
        for dll_path in vec_dll_paths:
            if os.path.exists(dll_path):
                try:
                    conn.load_extension(dll_path)
                    print(f"    ✅ 成功加载sqlite-vec: {dll_path}")
                    break
                except Exception as e:
                    print(f"    ⚠️  加载扩展失败: {e}")
                    continue
        
        conn.enable_load_extension(False)
    except Exception as e:
        print(f"    ⚠️  无法启用扩展: {e}")
    
    cursor = conn.cursor()
    
    # 创建文档图片表
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS doc_images (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            doc_id TEXT NOT NULL,
            image_hash TEXT NOT NULL,
            filename TEXT NOT NULL,
            file_path TEXT NOT NULL,
            description TEXT DEFAULT '',
            embedding TEXT,
            description_model TEXT DEFAULT '',
            embedding_model TEXT DEFAULT '',
            file_size INTEGER DEFAULT 0,
            created_at TEXT DEFAULT (datetime('now','localtime'))
        )
    """)
    
    # 创建向量表（使用与主系统相同的配置）
    try:
        # 先检查是否已存在
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='doc_image_vectors'")
        if not cursor.fetchone():
            cursor.execute("""
                CREATE VIRTUAL TABLE doc_image_vectors USING vec0(
                    id TEXT PRIMARY KEY,
                    vector FLOAT[1536]
                )
            """)
    except Exception as e:
        print(f"    ⚠️  向量表创建失败: {e}")
        print(f"    ℹ️  将使用模糊搜索代替向量搜索")
    
    conn.commit()
    conn.close()


def save_image_to_db(doc_id: str, image_info: Dict, description: str, 
                     desc_model: str, embedding: List[float], embed_model: str):
    """保存图片信息到数据库"""
    import sqlite3
    
    init_document_image_table()
    
    conn = sqlite3.connect(str(IMAGE_VECTORS_DB))
    cursor = conn.cursor()
    
    # 插入图片记录
    cursor.execute("""
        INSERT INTO doc_images 
        (doc_id, image_hash, filename, file_path, description, embedding, description_model, embedding_model, file_size)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
    """, (
        doc_id,
        image_info['hash'],
        image_info['filename'],
        image_info['filepath'],
        description,
        json.dumps(embedding) if embedding else None,
        desc_model,
        embed_model,
        image_info.get('size', 0)
    ))
    
    # 插入向量
    if embedding:
        vec_id = f"{doc_id}_{image_info['hash']}"
        try:
            cursor.execute(
                "INSERT OR REPLACE INTO doc_image_vectors (id, vector) VALUES (?, ?)",
                (vec_id, json.dumps(embedding))
            )
        except Exception as e:
            print(f"    ⚠️  向量存储失败: {e}")
    
    conn.commit()
    conn.close()


def get_images_by_doc(doc_id: str) -> List[Dict]:
    """获取文档的所有图片"""
    import sqlite3
    
    conn = sqlite3.connect(str(IMAGE_VECTORS_DB))
    cursor = conn.cursor()
    
    cursor.execute("""
        SELECT id, doc_id, image_hash, filename, file_path, description, 
               description_model, embedding_model, file_size, created_at
        FROM doc_images
        WHERE doc_id = ?
        ORDER BY created_at
    """, (doc_id,))
    
    rows = cursor.fetchall()
    images = []
    
    for row in rows:
        images.append({
            'id': row[0],
            'doc_id': row[1],
            'image_hash': row[2],
            'filename': row[3],
            'file_path': row[4],
            'description': row[5],
            'description_model': row[6],
            'embedding_model': row[7],
            'file_size': row[8],
            'created_at': row[9]
        })
    
    conn.close()
    return images


def search_images(query: str, limit: int = 5) -> List[Dict]:
    """搜索相关图片"""
    import sqlite3
    import json
    
    conn = sqlite3.connect(str(IMAGE_VECTORS_DB))
    cursor = conn.cursor()
    
    results = []
    
    # 模糊搜索（始终执行）
    try:
        cursor.execute("""
            SELECT id, doc_id, filename, description, file_path
            FROM doc_images
            WHERE description LIKE ? OR filename LIKE ?
            ORDER BY created_at DESC
            LIMIT ?
        """, (f'%{query}%', f'%{query}%', limit))
        rows = cursor.fetchall()
        for row in rows:
            results.append({
                'id': row[0],
                'doc_id': row[1],
                'filename': row[2],
                'description': row[3],
                'file_path': row[4],
                'similarity': 0.8
            })
        print(f"    ℹ️  模糊搜索返回 {len(results)} 条结果")
    except Exception as e:
        print(f"    ⚠️  模糊搜索失败: {e}")
    
    conn.close()
    return results


# ==================== 主处理函数 ====================

def process_document_images(filepath: str, doc_id: str = None) -> Dict:
    """
    处理文档，提取图片并生成描述和向量
    返回处理结果
    """
    import uuid
    
    # 生成文档ID
    if not doc_id:
        doc_id = f"doc_{uuid.uuid4().hex[:8]}"
    
    # 扩展名
    ext = Path(filepath).suffix.lower()
    
    print(f"\n📄 处理文档: {Path(filepath).name}")
    print(f"   文档ID: {doc_id}")
    
    # 提取图片
    print(f"\n🔍 提取图片...")
    if ext == '.docx':
        images = extract_images_from_docx(filepath)
    elif ext == '.pptx':
        images = extract_images_from_pptx(filepath)
    else:
        print(f"  ❌ 不支持的文件类型: {ext}")
        return {'success': False, 'error': f'不支持的文件类型: {ext}'}
    
    print(f"   找到 {len(images)} 张图片")
    
    if not images:
        return {'success': False, 'error': '未找到图片'}
    
    # 处理每张图片
    processed = []
    for i, img in enumerate(images, 1):
        print(f"\n   [{i}/{len(images)}] 处理: {img['filename']}")
        
        # 生成描述
        print(f"      生成描述...")
        description, desc_model = describe_image(img['filepath'])
        print(f"      模型: {desc_model}, 长度: {len(description)} 字")
        
        # 生成向量
        print(f"      生成向量...")
        embedding, embed_model = generate_embedding(description)
        print(f"      模型: {embed_model}, 维度: {len(embedding)}")
        
        # 保存到数据库
        save_image_to_db(doc_id, img, description, desc_model, embedding, embed_model)
        
        processed.append({
            'filename': img['filename'],
            'description': description[:100] + '...',
            'model': desc_model
        })
    
    print(f"\n✅ 处理完成！共处理 {len(processed)} 张图片")
    
    return {
        'success': True,
        'doc_id': doc_id,
        'image_count': len(processed),
        'images': processed
    }


# ==================== 统计 ====================

def get_stats() -> Dict:
    """获取统计信息"""
    import sqlite3
    
    conn = sqlite3.connect(str(IMAGE_VECTORS_DB))
    cursor = conn.cursor()
    
    try:
        cursor.execute("SELECT COUNT(*) FROM doc_images")
        total_images = cursor.fetchone()[0]
        
        cursor.execute("SELECT COUNT(DISTINCT doc_id) FROM doc_images")
        total_docs = cursor.fetchone()[0]
        
        result = {
            'total_images': total_images,
            'total_docs': total_docs
        }
    except:
        result = {
            'total_images': 0,
            'total_docs': 0
        }
    
    conn.close()
    return result


# ==================== 图片文件服务 ====================

def get_image_file(image_id: int):
    """获取图片文件"""
    import sqlite3
    
    conn = sqlite3.connect(str(IMAGE_VECTORS_DB))
    cursor = conn.cursor()
    
    cursor.execute("SELECT file_path FROM doc_images WHERE id = ?", (image_id,))
    row = cursor.fetchone()
    conn.close()
    
    if not row:
        return None
    
    file_path = row[0]
    if os.path.exists(file_path):
        from fastapi.responses import FileResponse
        return FileResponse(file_path)
    
    return None


if __name__ == '__main__':
    # 测试
    import sys
    
    if len(sys.argv) > 1:
        filepath = sys.argv[1]
        result = process_document_images(filepath)
        print(f"\n结果: {json.dumps(result, ensure_ascii=False, indent=2)}")
    else:
        print("用法: python knowledge_image_processor.py <文档路径>")
